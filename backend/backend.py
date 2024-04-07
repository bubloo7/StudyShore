from flask import Flask, request, send_file
from flask_cors import CORS
import os
from docx import Document
from pypdf import PdfReader
from pptx import Presentation
import random
import string
import json
import sys


from moviepy.editor import VideoFileClip
from quiz_generator import export_quiz
from flash_card_generator import export_flashcards
from summary_generator import export_summary
from gemini_call import prompt_everyting
from speech_to_text import get_audio

app = Flask(__name__)
CORS(app)


def generate_quiz(flashcards):
    quiz = []
    for card in flashcards:
        if random.choice([True, False]):
            question = {"question": card[1], "possible_answers": [], "index": -1}
            incorrect_answers = []
            for other_card in flashcards:
                if other_card != card:
                    incorrect_answers.append(other_card[0])
            question["possible_answers"] = random.sample(incorrect_answers, 3)
            question["index"] = random.randint(0, 3)
            question["possible_answers"].insert(question["index"], card[0])
        else:
            question = {"question": card[0], "possible_answers": [], "index": -1}
            incorrect_answers = []
            for other_card in flashcards:
                if other_card != card:
                    incorrect_answers.append(other_card[1])
            question["possible_answers"] = random.sample(incorrect_answers, 3)
            question["index"] = random.randint(0, 3)
            question["possible_answers"].insert(question["index"], card[1])
        quiz.append(question)
    return quiz


def generate_id():
    ids = set()
    for file in os.listdir(os.getcwd() + "/database"):
        if file.endswith(".json"):
            ids.add(file)

    letters = string.ascii_letters + string.digits
    id = "".join(random.choice(letters) for _ in range(5))
    while id in ids:
        id = "".join(random.choice(letters) for _ in range(5))
    return id


def example_response():
    flash_cards = [
        [
            "Photosynthesis",
            "The process by which green plants and some other organisms use sunlight to synthesize foods with the help of chlorophyll.",
        ],
        [
            "Mitochondria",
            "Organelles that generate most of the chemical energy needed to power the biochemical reactions of cells.",
        ],
        [
            "Newton's Laws",
            "Three fundamental principles of classical mechanics proposed by Sir Isaac Newton.",
        ],
        [
            "H2O",
            "Chemical formula for water, composed of two hydrogen atoms bonded to one oxygen atom.",
        ],
        [
            "Civil Rights Movement",
            "A struggle for social justice that took place mainly during the 1950s and 1960s for African Americans to gain equal rights under the law in the United States.",
        ],
        [
            "Palindrome",
            "A word, phrase, number, or other sequence of characters that reads the same forward and backward.",
        ],
        [
            "The Great Depression",
            "A severe worldwide economic depression that took place mostly during the 1930s.",
        ],
        [
            "E=mc^2",
            "Albert Einstein's famous equation, which expresses the relationship between energy (E), mass (m), and the speed of light (c) squared.",
        ],
        [
            "Renaissance",
            "A period in European history marking the transition from the Middle Ages to modernity and covering the 14th to 17th centuries.",
        ],
        [
            "Cell Division",
            "The process by which a parent cell divides into two or more daughter cells.",
        ],
    ]
    summary = """Lorem ipsum dolor sit amet, consectetur adipiscing elit, sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. Pharetra vel turpis nunc eget lorem dolor sed. Diam vulputate ut pharetra sit amet. Amet dictum sit amet justo. Nibh praesent tristique magna sit amet purus gravida quis. Nunc id cursus metus aliquam eleifend mi in nulla posuere. Sollicitudin aliquam ultrices sagittis orci. Egestas erat imperdiet sed euismod nisi porta lorem mollis. In fermentum posuere urna nec tincidunt praesent semper feugiat nibh. Vel pharetra vel turpis nunc eget lorem. Suspendisse sed nisi lacus sed viverra tellus in hac habitasse. Adipiscing enim eu turpis egestas. Facilisis gravida neque convallis a cras semper. Quam quisque id diam vel quam elementum pulvinar etiam.

Varius duis at consectetur lorem donec. Mi bibendum neque egestas congue quisque egestas diam in. Id cursus metus aliquam eleifend mi. Lacus laoreet non curabitur gravida arcu ac tortor dignissim convallis. Risus commodo viverra maecenas accumsan lacus vel facilisis volutpat est. Massa id neque aliquam vestibulum morbi blandit cursus. Commodo nulla facilisi nullam vehicula. Nec ullamcorper sit amet risus nullam eget. Pretium lectus quam id leo in vitae turpis massa sed. Elementum curabitur vitae nunc sed velit dignissim.

Posuere sollicitudin aliquam ultrices sagittis orci a scelerisque purus semper. Eget arcu dictum varius duis at consectetur. Tellus at urna condimentum mattis. At lectus urna duis convallis convallis. Sollicitudin nibh sit amet commodo. Quam viverra orci sagittis eu volutpat. Fermentum leo vel orci porta non pulvinar neque laoreet suspendisse. Amet est placerat in egestas erat imperdiet sed. Quisque non tellus orci ac auctor augue mauris augue neque. Feugiat sed lectus vestibulum mattis ullamcorper velit sed ullamcorper. Cursus sit amet dictum sit. In vitae turpis massa sed elementum tempus. Iaculis nunc sed augue lacus. Cras fermentum odio eu feugiat pretium.

Aliquam nulla facilisi cras fermentum. Nunc mattis enim ut tellus elementum. Tortor dignissim convallis aenean et tortor. Augue eget arcu dictum varius duis at consectetur lorem donec. Nunc id cursus metus aliquam eleifend mi in nulla posuere. Pellentesque massa placerat duis ultricies lacus sed turpis tincidunt id. Risus quis varius quam quisque. Commodo elit at imperdiet dui accumsan sit amet. Nunc sed augue lacus viverra. Commodo nulla facilisi nullam vehicula ipsum a arcu. Quis enim lobortis scelerisque fermentum dui faucibus in ornare.

Sodales neque sodales ut etiam sit amet nisl. Fames ac turpis egestas maecenas pharetra convallis posuere morbi leo. Sagittis id consectetur purus ut faucibus pulvinar. Et pharetra pharetra massa massa. Posuere sollicitudin aliquam ultrices sagittis orci. Nisl nisi scelerisque eu ultrices vitae auctor eu augue. Ut tellus elementum sagittis vitae et leo duis ut diam. Convallis posuere morbi leo urna molestie at. Urna et pharetra pharetra massa. Neque sodales ut etiam sit. Risus viverra adipiscing at in. Dui sapien eget mi proin sed libero enim sed. Purus sit amet volutpat consequat mauris nunc. Nunc non blandit massa enim nec dui nunc."""

    quiz = generate_quiz(flash_cards)

    id = generate_id()

    letters = string.ascii_letters + string.digits
    title = "Title " + "".join(random.choice(letters) for _ in range(8))
    db_json = {
        "summary": summary,
        "flash_cards": flash_cards,
        "quiz": quiz,
        "title": title,
    }

    with open(os.getcwd() + "/database/" + id + ".json", "w") as f:
        json.dump(db_json, f)

    return {"id": id}


def handle_pdf(pdf):
    reader = PdfReader(os.getcwd() + "/file.pdf")
    number_of_pages = len(reader.pages)
    s = ""
    for i in range(number_of_pages):
        page = reader.pages[i]
        text = page.extract_text()
        s += text

    return s


def handle_txt(txt):
    s = open("file.txt").read()
    print(s, "save me god")
    return s


def handle_docx(docx):
    d = Document(os.getcwd() + "/file.docx")
    s = ""
    for paragraph in d.iter_inner_content():
        for line in paragraph.iter_inner_content():
            s += line.text
        s += "\n"

    return s


def handle_mp3(mp3):
    # TODO: the mp3 file that was uploaded should be saved as file.mp3. Convert the mp3 to string/transcript, then return the string/transcript. That is all
    s = get_audio("file.mp3")
    print("handle mp3 string", s)
    return s


def handle_mp4(mp4):
    # TODO Converts the mp4 file to mp3 and returns what that returns. I think it should work but it is worth testing
    video = VideoFileClip("file.mp4")
    audio = video.audio
    audio.write_audiofile("file.mp3")
    audio.close()
    video.close()
    return handle_mp3("file.mp3")


def handle_pptx(pptx):
    p = Presentation("file.pptx")
    s = ""
    # sorry about the try/except stuff. its annoying bc certain things dont have texts/shapes and make it crash
    for slide in p.slides:
        try:
            for shape in slide.shapes:
                try:
                    s += shape.text
                except Exception:
                    pass
            s += "\n"
        except Exception:
            pass

    return s


@app.route("/")
def hello_world():
    return "Hello, World!"


@app.route("/upload", methods=["POST"])
def upload():
    print("in /upload")
    file = request.files["file"]
    name = file.filename
    extension = name.split(".")[-1]
    file.save(dst=os.getcwd() + "/file." + extension)
    file.close()

    if extension == "pdf":
        s = handle_pdf(file)
    elif extension == "txt":
        s = handle_txt(file)
    elif extension == "docx":
        s = handle_docx(file)
    elif extension == "mp3":
        s = handle_mp3(file)
    elif extension == "mp4":
        s = handle_mp4(file)
    elif extension == "pptx":
        s = handle_pptx(file)

    # print("parsed s:", s)
    response = prompt_everyting(s)
    # print("finished prompting")
    # print(response)

    quiz2 = generate_quiz(response["flash_cards"])
    response["quiz"] += quiz2

    random.shuffle(response["quiz"])

    id = generate_id()
    response["id"] = id
    with open(os.getcwd() + "/database/" + id + ".json", "w") as f:
        json.dump(response, f)

    return response


@app.route("/fetch_id", methods=["POST"])
def fetch_id():
    id = request.json.get("id")
    try:
        with open(os.getcwd() + "/database/" + id + ".json", "r") as f:
            data = json.load(f)

        print(data)
        return data
    except:
        return {"summary": 404}


@app.route("/recent", methods=["GET"])
def recent():
    print("entering recent")

    output = {"recent": []}
    i = 0
    for file in os.listdir(os.getcwd() + "/database"):
        if i == 10:
            break
        if file.endswith(".json"):
            with open(os.getcwd() + "/database/" + file, "r") as f:
                data = json.load(f)
                id = file.split(".json")[0]
                title = data["title"]
                output["recent"].append({"id": id, "title": title})

        i += 1
    print("leaving recent")
    return output


@app.route("/export", methods=["POST"])
def export():
    selected = request.json.get("selected")
    data = request.json.get("data")
    if selected == 0:
        data = data["summary"]
        export_summary(data, "Summary.docx")
        return send_file("Summary.docx", as_attachment=True)
    elif selected == 1:
        data = data["flash_cards"]
        export_flashcards(data, "Flashcards.docx")
        return send_file("Flashcards.docx", as_attachment=True)

    else:
        data = data["quiz"]

        export_quiz(data, "Quiz.docx")
        return send_file("Quiz.docx", as_attachment=True)


if __name__ == "__main__":
    app.run(debug=True, port=5000)
